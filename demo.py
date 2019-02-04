import os
import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def makedir(path):
    try:
        os.makedirs(path)
    except OSError as exc:
        if os.path.isdir(path):
            pass
        else:
            raise

def make_pptx(book,chapter,verses):

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    mycolor = RGBColor(0xef, 0xef, 0xff)

    for i, verse in enumerate(verses):
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture('./background.jpg',Inches(0),Inches(0),Inches(10),Inches(7.5))

        txBox = slide.shapes.add_textbox(Inches(0.5),Inches(0.5),Inches(9),Inches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = book+' '+chapter+':'+str(i+1)
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = mycolor

        p = tf.add_paragraph()
        p.text = ''
        p.font.size = Pt(10)

        p = tf.add_paragraph()
        p.text = '    '+verse
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = mycolor

    print(book, chapter, len(verses))
    
    makedir('./outputs/'+book)
    prs.save('./outputs/'+book+'/'+book+chapter+'.pptx')

if __name__ == '__main__':
    csvfile = open('./samples.csv', 'r', newline='', encoding='utf-8')
    fields = ['book','chapter','verse','content']
    reader = csv.DictReader(csvfile, fieldnames=fields)

    verses = []
    for line in reader:
        if '\ufeff' in line['book']: # there're some strange characters at the beginning of this csv file.
            line['book'] = line['book'][-1:]
            current_book = line['book']
            current_chapter = '1'
        if current_book != line['book'] or current_chapter != line['chapter']:
            make_pptx(current_book,current_chapter,verses)
            current_book = line['book']
            current_chapter = line['chapter']
            verses = []
        verses.append(line['content'])
    make_pptx(line['book'],current_chapter,verses)

    csvfile.close()